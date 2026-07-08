#!/usr/bin/env python3
"""Build the FirstPass interview presentation (deck/FirstPass_PortSwigger.pptx).

PortSwigger brand: white background, near-black text, orange (#FF6633) accents,
GREEN/AMBER/RED triage motif. Interview-presentation arc (how I think -> the
scenario -> the solution -> results -> value). No accent bars / edge stripes /
title underlines (AI-slide tells). No em-dashes (house style). 15 slides.
Reproducible: rerun to regenerate.
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
ASSETS = ROOT / "assets"

# ---- palette (light) ---------------------------------------------------------
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
SOFT    = RGBColor(0xF5, 0xF6, 0xF8)
SOFT2   = RGBColor(0xEC, 0xEF, 0xF2)
LINE    = RGBColor(0xE1, 0xE5, 0xEA)
INK     = RGBColor(0x15, 0x17, 0x1C)
INK2    = RGBColor(0x39, 0x41, 0x50)
DIM     = RGBColor(0x5A, 0x64, 0x72)
FAINT   = RGBColor(0x9A, 0xA0, 0xAA)
ACC     = RGBColor(0xFF, 0x66, 0x33)   # PortSwigger orange
ACC2    = RGBColor(0xE1, 0x4E, 0x12)   # darker orange for text on white
GREEN   = RGBColor(0x12, 0xA1, 0x50)
AMBER   = RGBColor(0xC0, 0x7C, 0x00)
RED     = RGBColor(0xDE, 0x32, 0x26)
GREEN_T = RGBColor(0xE7, 0xF6, 0xEE)
AMBER_T = RGBColor(0xFB, 0xF2, 0xDC)
RED_T   = RGBColor(0xFB, 0xE9, 0xE7)
ORANGE_T= RGBColor(0xFC, 0xEB, 0xE3)
BLUE    = RGBColor(0x2F, 0x6F, 0xB0)
BLUE_T  = RGBColor(0xE8, 0xF1, 0xFB)
CHARC   = RGBColor(0x15, 0x17, 0x1C)

FONT = "Arial"
SW, SH = 13.333, 7.5
M = 0.85
_num = 0

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


# ---- helpers -----------------------------------------------------------------
def _noshadow(sh):
    try:
        sh.shadow.inherit = False
    except Exception:
        pass


def slide(bg=WHITE, footer=True):
    global _num
    _num += 1
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill
    f.solid(); f.fore_color.rgb = bg
    if footer:
        fc = FAINT if bg == WHITE else RGBColor(0x6B, 0x72, 0x80)
        text(s, "FirstPass  ·  AI Pioneer  ·  Abdelrhman Rayis", M, SH - 0.46, 8, 0.3, 9.5, fc)
        text(s, f"{_num:02d}", SW - 1.0, SH - 0.46, 0.5, 0.3, 9.5, fc, align=PP_ALIGN.RIGHT)
    return s


def text(s, txt, l, t, w, h, size, color=INK, bold=False, italic=False,
         align=PP_ALIGN.LEFT, spacing=1.04, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = spacing
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.name = font; r.font.color.rgb = color
    return tb


def multi(s, lines, l, t, w, h, size, spacing=1.18, align=PP_ALIGN.LEFT, space_after=6):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, item in enumerate(lines):
        content, c, b = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = spacing; p.space_after = Pt(space_after)
        r = p.add_run(); r.text = content
        r.font.size = Pt(size); r.font.bold = b; r.font.name = FONT; r.font.color.rgb = c
    return tb


def card(s, l, t, w, h, fill=SOFT, line=LINE, radius=0.045):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is not None:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    _noshadow(sh)
    return sh


def circle(s, label, l, t, d=0.52, fill=ACC, tcolor=WHITE, size=15):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = fill; c.line.fill.background(); _noshadow(c)
    tf = c.text_frame; tf.word_wrap = False; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(size); r.font.bold = True; r.font.name = FONT; r.font.color.rgb = tcolor
    return c


def dot(s, color, l, t, d=0.2):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    o.fill.solid(); o.fill.fore_color.rgb = color; o.line.fill.background(); _noshadow(o)
    return o


def pill(s, txt, l, t, w, h, fill, txt_color, size=12, bold=True):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    b.adjustments[0] = 0.5
    b.fill.solid(); b.fill.fore_color.rgb = fill; b.line.fill.background(); _noshadow(b)
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONT; r.font.color.rgb = txt_color
    return b


def node(s, txt, l, t, w, h, fill=SOFT, tcolor=INK, size=11.5, line=LINE, bold=True):
    c = card(s, l, t, w, h, fill=fill, line=line, radius=0.1)
    tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Emu(9144); tf.margin_right = Emu(9144)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = FONT; r.font.color.rgb = tcolor
    return c


def title(s, kicker, ttl, tcolor=INK):
    text(s, kicker.upper(), M, 0.52, 11.6, 0.35, 12.5, ACC2, bold=True)
    text(s, ttl, M, 0.92, SW - 2 * M, 1.2, 31, tcolor, bold=True)


def arrow(s, l, t, color=ACC, size=17):
    text(s, "›", l, t, 0.3, 0.3, size, color, bold=True, align=PP_ALIGN.CENTER)


def image(s, path, l, t, height, border=True):
    pic = s.shapes.add_picture(str(path), Inches(l), Inches(t), height=Inches(height))
    if border:
        pic.line.color.rgb = LINE; pic.line.width = Pt(1.25)
    _noshadow(pic)
    return pic


# =============================================================================
# 1. COVER
# =============================================================================
s = slide(footer=False)
dot(s, GREEN, M, 1.62, 0.2); dot(s, AMBER, M + 0.32, 1.62, 0.2); dot(s, RED, M + 0.64, 1.62, 0.2)
text(s, "FirstPass", M - 0.04, 1.95, 10, 1.4, 82, INK, bold=True)
text(s, "Clearing the contract-review queue, and how I think.", M, 3.45, 8.4, 0.6, 23, INK, bold=True)
text(s, "A working prototype, a live console, and measured results,", M, 4.12, 11, 0.4, 15.5, DIM)
text(s, "built to run entirely on your own servers.", M, 4.42, 11, 0.4, 15.5, DIM)
text(s, "Abdelrhman Rayis", M, 5.7, 8, 0.4, 16, INK, bold=True)
text(s, "AI Pioneer  ·  work-sample presentation", M, 6.08, 8, 0.4, 12.5, DIM)
# right column
text(s, "RISK TRIAGE", SW - 3.55, 2.0, 2.8, 0.4, 11.5, FAINT, bold=True, align=PP_ALIGN.RIGHT)
pill(s, "RED    escalate", SW - 3.55, 2.42, 2.8, 0.5, RED_T, RED, size=12.5)
pill(s, "AMBER    confirm", SW - 3.55, 3.02, 2.8, 0.5, AMBER_T, AMBER, size=12.5)
pill(s, "GREEN    auto-clear", SW - 3.55, 3.62, 2.8, 0.5, GREEN_T, GREEN, size=12.5)
text(s, "\U0001F512  on-prem  ·  local AI  ·  no cloud", SW - 3.9, 4.35, 3.15, 0.4, 11.5, GREEN,
     bold=True, align=PP_ALIGN.RIGHT)

# =============================================================================
# 2. AGENDA
# =============================================================================
s = slide()
title(s, "What I'll cover", "Five short steps, the way I'd tackle any new problem.")
agenda = [
    ("How I approach a problem I have not seen before", "First principles, five whys, systems thinking, 80/20."),
    ("The scenario, and its root cause", "Your contract queue, and why it really backs up."),
    ("FirstPass: the solution, running live", "An agentic triage layer, in a browser, on your server."),
    ("Does it actually help?", "Measured properly: safety, accuracy, hours saved."),
    ("Private by design, and how I would add value", "Local AI, the standards that matter, and my fit."),
]
y = 2.2
for i, (h, sub) in enumerate(agenda, 1):
    circle(s, str(i), M, y, 0.5, fill=ACC, size=15)
    text(s, h, M + 0.8, y - 0.04, 10.6, 0.4, 16.5, INK, bold=True)
    text(s, sub, M + 0.8, y + 0.36, 10.6, 0.35, 12.5, DIM)
    y += 0.92

# =============================================================================
# 3. THE OVERALL WORKFLOW  (agile delivery)
# =============================================================================
s = slide()
title(s, "The plan", "Value from day one: a skill that grows into an agentic loop.")
stages = [
    ("DAY 1", "Skill", "Deploy the contract-triage skill, after understanding the business."),
    ("DAY 2", "Agent memory", "The playbook becomes the agent's memory of standard positions."),
    ("WEEK 1", "System", "The console, routing and evals grow around it."),
    ("WEEK 2", "Agentic loop", "The full loop, with the learning flywheel."),
]
cw = 2.52; gp = 0.49; x = M
for i, (when, name, desc) in enumerate(stages):
    pill(s, when, x + cw / 2 - 0.62, 2.18, 1.24, 0.4, ORANGE_T, ACC2, size=11)
    card(s, x, 2.72, cw, 1.68, fill=SOFT)
    circle(s, str(i + 1), x + 0.26, 2.94, 0.42, fill=ACC, size=13)
    text(s, name, x + 0.8, 2.99, cw - 0.95, 0.4, 15, INK, bold=True)
    text(s, desc, x + 0.28, 3.52, cw - 0.52, 0.85, 11.5, DIM, spacing=1.16)
    if i < len(stages) - 1:
        arrow(s, x + cw + 0.1, 3.36, size=20)
    x += cw + gp
card(s, M, 4.62, SW - 2 * M, 0.55, fill=ORANGE_T, line=None, radius=0.14)
text(s, "↺   Feedback loop: gather feedback and improve every two weeks. The queue gets shorter each cycle.",
     M + 0.35, 4.74, SW - 2 * M - 0.6, 0.35, 12.5, ACC2, bold=True)
card(s, M, 5.42, SW - 2 * M, 1.05, fill=SOFT)
text(s, "WHERE IT RUNS", M + 0.34, 5.58, 4, 0.3, 11, FAINT, bold=True)
pill(s, "Cloud, fast: synthetic test data, no GDPR issue", M + 0.34, 5.92, 3.85, 0.44, SOFT2, DIM, size=11)
arrow(s, M + 4.28, 6.04, size=16)
pill(s, "\U0001F512  Local Hermes, on your server: real data", M + 4.62, 5.92, 4.15, 0.44, GREEN_T, GREEN, size=11)
text(s, "When real client data is involved, reasoning moves in-house. Value fast, privacy kept.",
     M + 9.0, 5.98, SW - M - (M + 9.0), 0.6, 11, INK2, italic=True, spacing=1.14)

# =============================================================================
# 4. HOW I APPROACH A PROBLEM  (mental models)
# =============================================================================
s = slide()
title(s, "How I think", "I don't jump to a solution. I find the root of the problem first.")
text(s, "Four mental models I reach for before writing a line of code:", M, 2.02, 11, 0.4, 14, DIM)
models = [
    ("First principles", "Strip the task to its atoms. What is really being decided here?"),
    ("Five whys", "Keep asking why until the symptom gives up its root cause."),
    ("Systems thinking", "Look at the whole flow. Where is the constraint, and the leverage?"),
    ("80 / 20", "Find the small share of cases that carries most of the value and risk."),
]
cw = (SW - 2 * M - 3 * 0.3) / 4
for i, (name, desc) in enumerate(models):
    x = M + i * (cw + 0.3)
    card(s, x, 2.6, cw, 2.85, fill=SOFT)
    dot(s, ACC, x + 0.32, 2.92, 0.26)
    text(s, name, x + 0.3, 3.42, cw - 0.5, 0.4, 15.5, INK, bold=True)
    text(s, desc, x + 0.3, 3.86, cw - 0.54, 1.4, 12.5, DIM, spacing=1.2)
text(s, "The point: a creative solution is worthless if it solves the wrong problem. Models keep me honest.",
     M, 5.75, SW - 2 * M, 0.4, 13.5, INK2, italic=True)

# =============================================================================
# 4. THE SCENARIO (problem)
# =============================================================================
s = slide()
title(s, "The scenario", "The case in front of me: a queue longer than the day.")
card(s, M, 2.02, SW - 2 * M, 1.62, fill=SOFT)
text(s, "“Most are 80% boilerplate. Where we get redlines, they are usually edits we have seen "
        "before. Right now, every single contract still gets read by a person, and the queue is "
        "usually longer than the day. At end of month or quarter it can get overwhelming.”",
     M + 0.35, 2.24, SW - 2 * M - 0.7, 1.1, 15, INK, italic=True, spacing=1.18)
text(s, "PortSwigger, the task brief", M + 0.35, 3.32, 6, 0.3, 11, FAINT)
tiles = [("80%", "of a typical contract is boilerplate, agreed a hundred times before", AMBER),
         ("100%", "still get a full cold read by a scarce, expensive human", RED),
         ("EoQ", "sales spikes turn the queue into a bottleneck on revenue", ACC2)]
tw = (SW - 2 * M - 2 * 0.3) / 3
for i, (big, lbl, col) in enumerate(tiles):
    x = M + i * (tw + 0.3)
    card(s, x, 3.95, tw, 1.62, fill=SOFT)
    text(s, big, x + 0.3, 4.14, tw - 0.5, 0.7, 34, col, bold=True)
    text(s, lbl, x + 0.3, 4.86, tw - 0.55, 0.6, 12.5, DIM, spacing=1.1)
text(s, "The cost is not just time. A slow queue slows every deal sitting behind it.",
     M, 5.82, 11, 0.4, 13.5, INK2, italic=True)

# =============================================================================
# 5. FIVE WHYS -> ROOT CAUSE
# =============================================================================
s = slide()
title(s, "Five whys", "Drilling past the symptom to the root of the issue.")
whys = [
    ("Why is the queue longer than the day?", "Every contract gets a full human read."),
    ("Why does each one need a full read?", "It is read cold, start to finish, every time."),
    ("Why is it read cold?", "There is no way to tell the safe ones from the risky ones up front."),
    ("Why can't we tell them apart?", "The team's standard positions live in people's heads, not in a checkable form."),
    ("Why is that the case?", "The process was built around people, with no triage layer in front of it."),
]
y = 1.95
for i, (q, a) in enumerate(whys, 1):
    circle(s, str(i), M, y, 0.42, fill=INK, size=13)
    text(s, q, M + 0.66, y - 0.02, 5.2, 0.5, 13.5, INK, bold=True, spacing=1.05)
    arrow(s, M + 5.95, y + 0.02, color=ACC, size=15)
    text(s, a, M + 6.4, y - 0.02, SW - M - (M + 6.4), 0.5, 13.5, DIM, spacing=1.05)
    y += 0.72
card(s, M, 5.62, SW - 2 * M, 1.05, fill=ORANGE_T, line=None)
text(s, "ROOT CAUSE", M + 0.35, 5.78, 3, 0.3, 11, ACC2, bold=True)
text(s, "A lawyer's scarce attention is spent identically on the safe 80% and the risky 20%. "
        "There is no risk-based triage. That is what I need to fix, not the reading speed.",
     M + 0.35, 6.08, SW - 2 * M - 0.7, 0.5, 14, INK, bold=True, spacing=1.12)

# =============================================================================
# 6. FIRST PRINCIPLES: RISK GATE
# =============================================================================
s = slide()
title(s, "First principles", "Reframed: contract review is a risk gate, not a reading task.")
text(s, "Strip it to the atoms and every review is the same five decisions:", M, 1.98, 11, 0.4, 14, DIM)
steps = ["Read", "Classify", "Compare\nto standard", "Spot the\ndeviation", "Accept /\nredline /\nescalate"]
w = 2.02; gap = 0.34; x = M
for i, st in enumerate(steps):
    node(s, st, x, 2.5, w, 1.15, fill=SOFT, size=13)
    circle(s, str(i + 1), x + 0.14, 2.64, 0.34, fill=ACC, size=11)
    if i < len(steps) - 1:
        arrow(s, x + w + 0.03, 2.95, size=19)
    x += w + gap
card(s, M, 4.15, SW - 2 * M, 1.85, fill=SOFT)
text(s, "THE INSIGHT", M + 0.35, 4.35, 4, 0.35, 11.5, ACC2, bold=True)
multi(s, [
    ("Only two of those five steps need legal judgement. The rest is mechanical.", INK, True),
    ("80% of contracts are boilerplate, and the redlines that come up are edits seen before: a "
     "finite, reusable playbook. So the share needing net-new legal thinking is small, and it shrinks "
     "as the playbook fills. Yet today one scarce resource, a lawyer's attention, is spent identically "
     "on the safe 80% and the risky 20%.", INK2, False),
], M + 0.35, 4.68, SW - 2 * M - 0.7, 1.2, 14.5, spacing=1.18, space_after=6)

# =============================================================================
# 7. SYSTEMS THINKING + BURP
# =============================================================================
s = slide()
title(s, "Systems thinking", "The fix is structural: insert a triage layer.")
text(s, "TODAY", M, 2.0, 3, 0.3, 12, RED, bold=True)
for i, c in enumerate(["NDA", "Form", "NDA", "Form"]):
    node(s, c, M + i * 0.74, 2.36, 0.64, 0.52, fill=SOFT, size=10, tcolor=DIM)
arrow(s, M + 2.98, 2.46, size=16)
node(s, "one human queue", M + 3.28, 2.24, 1.75, 0.76, fill=RED, tcolor=WHITE, size=11, line=None)
text(s, "Every item pays the same full cost. Volume spikes have nowhere to go.",
     M + 5.25, 2.42, 6.2, 0.5, 12.5, INK2, italic=True)
text(s, "WITH FIRSTPASS", M, 3.4, 4, 0.3, 12, GREEN, bold=True)
node(s, "contracts", M, 3.78, 1.35, 0.85, fill=SOFT, size=11)
arrow(s, M + 1.46, 4.04, size=16)
node(s, "triage layer", M + 1.76, 3.78, 1.6, 0.85, fill=ACC, tcolor=WHITE, size=11, line=None)
pill(s, "GREEN  auto-clear", M + 3.62, 3.66, 2.95, 0.44, GREEN_T, GREEN, size=11.5)
pill(s, "AMBER  confirm", M + 3.62, 4.16, 2.95, 0.44, AMBER_T, AMBER, size=11.5)
pill(s, "RED  escalate", M + 3.62, 4.66, 2.95, 0.44, RED_T, RED, size=11.5)
card(s, M + 6.95, 3.55, SW - M - (M + 6.95), 1.6, fill=SOFT)
text(s, "You already do this", M + 7.25, 3.72, 4, 0.35, 13, ACC2, bold=True)
text(s, "Burp Suite triages security findings by severity so humans focus on the real threats. "
        "FirstPass is the same move, pointed at the contract queue.",
     M + 7.25, 4.08, SW - M - (M + 6.95) - 0.5, 1.0, 13, INK2, spacing=1.16)
text(s, "Leverage point: route attention by risk, so human judgement flows only to the 20% that needs it.",
     M, 5.55, SW - 2 * M, 0.4, 13.5, INK2, italic=True)

# =============================================================================
# PROCESS ANALYSIS  (business analyst view)
# =============================================================================
s = slide()
title(s, "Business analysis", "The process, its friction, and the right fix at each step.")
text(s, "Not everything should be AI. I picked the right tool for each point in the process.",
     M, 1.98, 11, 0.35, 13, DIM, italic=True)
rows = [
    ("Intake & parse", "Mixed formats, manual logging, no priority",
     "Auto-ingest and extract. This is n8n / RPA territory.", "AUTOMATION", BLUE, BLUE_T),
    ("Classify & pull terms", "Implicit, redone by hand every time",
     "The model reads and understands the document.", "AGENTIC", ACC2, ORANGE_T),
    ("Check the playbook", "Standard positions live in people's heads",
     "Retrieve the position, match the clause by meaning.", "AGENTIC", ACC2, ORANGE_T),
    ("Spot risk & red flags", "Fatigue, and the risky 20% slips through",
     "Rules for hard flags, the agent for paraphrase and novelty.", "HYBRID", AMBER, AMBER_T),
    ("Draft redline & route", "Re-writing edits seen before; first-in-first-out",
     "Draft from the playbook, score GREEN / AMBER / RED.", "AGENTIC", ACC2, ORANGE_T),
    ("Decide & capture", "Nothing is learned; novel issues re-solved each time",
     "The lawyer decides RED; the decision feeds the knowledge base.", "HUMAN + LOOP", GREEN, GREEN_T),
]
y = 2.5
for step, fric, fix, tag, col, tint in rows:
    text(s, step, M, y + 0.03, 2.7, 0.6, 13.5, INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    multi(s, [(fric, DIM, False), (fix, INK2, True)], M + 2.78, y, 5.35, 0.66, 11.5, spacing=1.12, space_after=2)
    pill(s, tag, M + 8.4, y + 0.14, 2.5, 0.4, tint, col, size=10.5)
    y += 0.7

# =============================================================================
# 8. THE AGENTIC LOOP
# =============================================================================
s = slide()
title(s, "The solution", "One agentic loop: contract in, a routed and redlined decision out.")
loop = ["Ingest", "Classify", "Retrieve\nplaybook", "Detect", "Score", "Draft\nredline", "Route"]
w = 1.5; gap = 0.16; x = M
for i, st in enumerate(loop):
    node(s, st, x, 2.35, w, 1.0, fill=SOFT, size=11)
    circle(s, str(i + 1), x + 0.12, 2.47, 0.3, fill=ACC, size=10)
    if i < len(loop) - 1:
        arrow(s, x + w - 0.01, 2.72, size=15)
    x += w + gap
card(s, M, 3.72, SW - 2 * M, 0.5, fill=ORANGE_T, line=None, radius=0.16)
text(s, "8   LEARN     the lawyer's decision is written back to the playbook, so next time it is routine",
     M + 0.3, 3.8, SW - 2 * M - 0.5, 0.35, 12.5, ACC2, bold=True)
oc = [("GREEN", GREEN, GREEN_T, "Auto-clear with a spot-check.\nNever auto-signed."),
      ("AMBER", AMBER, AMBER_T, "Known deviation. The redline is\nalready drafted. Lawyer confirms."),
      ("RED", RED, RED_T, "Novel or high risk. Escalate\nwith terms and issues pre-extracted.")]
w2 = (SW - 2 * M - 2 * 0.3) / 3
for i, (name, col, tint, desc) in enumerate(oc):
    xx = M + i * (w2 + 0.3)
    card(s, xx, 4.5, w2, 1.55, fill=SOFT)
    pill(s, name, xx + 0.3, 4.7, 1.55, 0.42, tint, col, size=13)
    text(s, desc, xx + 0.3, 5.26, w2 - 0.55, 0.9, 13, INK2, spacing=1.16)

# =============================================================================
# WHY AGENTIC, NOT n8n  (decision rationale)
# =============================================================================
s = slide()
title(s, "Why an agentic loop, not n8n", "Deterministic automation is the wrong tool for a judgement task.")
text(s, "n8n, Zapier and RPA are excellent when the logic can be fully specified up front. Contract "
        "wording never can: it is unstructured, endlessly paraphrased, and full of clauses no fixed "
        "rule has seen. That is a reasoning problem, not a plumbing one.",
     M, 1.95, SW - 2 * M, 0.7, 13.5, INK2, spacing=1.2)
caps = [
    ("Move data and call APIs (the plumbing)", "y", "p", "p"),
    ("Understand unstructured contract language", "n", "p", "y"),
    ("Handle paraphrased or novel clauses", "n", "n", "y"),
    ("Apply a risk policy and judgement", "n", "p", "y"),
    ("Improve from feedback over time", "n", "n", "y"),
    ("Explainable, with a human in the loop", "p", "y", "y"),
]
cols = [("n8n / RPA", 7.55), ("Rules only", 9.35), ("Agentic loop", 11.15)]
card(s, M, 2.75, SW - 2 * M, 3.05, fill=SOFT)
text(s, "WHAT THE JOB NEEDS", M + 0.35, 2.92, 5, 0.3, 11, FAINT, bold=True)
for name, cx in cols:
    text(s, name, cx - 0.9, 2.9, 1.8, 0.3, 11, FAINT, bold=True, align=PP_ALIGN.CENTER)
glyph = {"y": ("✓", GREEN), "p": ("~", AMBER), "n": ("–", FAINT)}
yy = 3.32
for cap, a, b, c in caps:
    text(s, cap, M + 0.35, yy, 6.0, 0.4, 12.5, INK, anchor=MSO_ANCHOR.MIDDLE)
    for val, (_, cx) in zip((a, b, c), cols):
        g, gc = glyph[val]
        text(s, g, cx - 0.9, yy - 0.02, 1.8, 0.4, 17, gc, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    yy += 0.4
card(s, M, 5.95, SW - 2 * M, 0.72, fill=ORANGE_T, line=None, radius=0.12)
text(s, "Not either / or: n8n moves the contracts, the agentic loop reviews them. Deterministic "
        "automation for the pipes, an agentic loop for the judgement, a human on the risky 20%.",
     M + 0.35, 6.06, SW - 2 * M - 0.7, 0.55, 13, ACC2, bold=True, spacing=1.14)

# =============================================================================
# 9. WEB CONSOLE (screenshot)
# =============================================================================
s = slide()
title(s, "It runs today", "A working console, not a mockup.")
image(s, ASSETS / "shot_dashboard.png", M, 1.95, height=4.8)
rx = M + 7.1
multi(s, [
    ("The queue a lawyer opens each morning.", INK, True),
    ("Sorted by risk. The safe 80% is already cleared, the redlines are drafted, and only the real "
     "risk is waiting for a human.", DIM, False),
], rx, 2.15, SW - M - rx, 1.6, 14, spacing=1.2, space_after=10)
card(s, rx, 4.15, SW - M - rx, 1.4, fill=GREEN_T, line=None)
text(s, "\U0001F512  Runs on your server", rx + 0.28, 4.34, 4, 0.35, 13, GREEN, bold=True)
text(s, "In-browser, local, offline. No contract text ever leaves the building.",
     rx + 0.28, 4.7, SW - M - rx - 0.55, 0.7, 12.5, INK2, spacing=1.16)

# =============================================================================
# 10. INSIDE A REVIEW (screenshot)
# =============================================================================
s = slide()
title(s, "Inside a review", "Every issue arrives with the redline already drafted.")
image(s, ASSETS / "shot_detail_amber.png", M, 1.95, height=4.95)
rx = M + 6.55
multi(s, [
    ("The lawyer confirms, they do not start from scratch.", INK, True),
    ("Key terms extracted. Each deviation shows the evidence, the plain-English reason, and the exact "
     "redline the team has used before.", DIM, False),
    ("Accept, edit or reject. One click.", INK, True),
], rx, 2.15, SW - M - rx, 2.4, 13.5, spacing=1.2, space_after=10)
text(s, "This is the artefact someone uses on Monday morning.", rx, 5.75, SW - M - rx, 0.6, 13, ACC2, italic=True, bold=True)

# =============================================================================
# 11. MEASURED, AND WORTH IT
# =============================================================================
s = slide()
title(s, "Does it help?", "Measured properly, on a labelled set, reproducibly.")
sd = [("0", "unsafe misses: nothing risky was auto-cleared", GREEN),
      ("100%", "triage verdict accuracy (10 of 10)", INK),
      ("93 / 93", "material-issue precision / recall, no LLM", AMBER),
      ("58%", "of human review minutes removed", ACC2)]
sw = (SW - 2 * M - 3 * 0.26) / 4
for i, (big, lbl, col) in enumerate(sd):
    x = M + i * (sw + 0.26)
    card(s, x, 2.05, sw, 1.7, fill=SOFT)
    text(s, big, x + 0.28, 2.26, sw - 0.4, 0.7, 31, col, bold=True)
    text(s, lbl, x + 0.28, 2.98, sw - 0.45, 0.7, 11.5, DIM, spacing=1.1)
card(s, M, 4.0, 6.15, 2.05, fill=ORANGE_T, line=None)
text(s, "EXTRAPOLATED TO 200 CONTRACTS / MONTH", M + 0.32, 4.2, 5.5, 0.35, 11.5, ACC2, bold=True)
text(s, "~73 hours", M + 0.32, 4.52, 5, 0.85, 40, INK, bold=True)
text(s, "of legal time returned every month, about two weeks of a reviewer, pointed back at real risk.",
     M + 0.32, 5.38, 5.6, 0.6, 12.5, INK2, spacing=1.12)
card(s, M + 6.4, 4.0, SW - M - (M + 6.4), 2.05, fill=SOFT)
text(s, "EXTERNAL BENCHMARK  ·  LAWGEEX 2018", M + 6.7, 4.2, 5, 0.35, 11.5, FAINT, bold=True)
multi(s, [
    ("94% vs 85%   AI beat 20 lawyers at NDA issue-spotting", GREEN, True),
    ("26 sec vs 92 min   average review time", GREEN, True),
    ("The category is proven. This applies it to your queue, privately.", DIM, False),
], M + 6.7, 4.6, SW - M - (M + 6.4) - 0.6, 1.3, 12.5, spacing=1.25, space_after=7)

# =============================================================================
# 12. KNOWLEDGE BASE + SKILLS THAT GROW
# =============================================================================
s = slide()
title(s, "It gets smarter", "A knowledge base and a skill set that grow with every case.")
cyc = [("A lawyer resolves\na novel case", RED, RED_T),
       ("It becomes a rule\nin the knowledge base", ACC, ORANGE_T),
       ("That case is now\nroutine, redline ready", AMBER, AMBER_T),
       ("Next month's queue\nis shorter", GREEN, GREEN_T)]
w = 2.55; gap = 0.42; x = M
for i, (t, col, tint) in enumerate(cyc):
    node(s, t, x, 2.25, w, 1.15, fill=tint, tcolor=INK, size=12.5, line=None)
    if i < len(cyc) - 1:
        arrow(s, x + w + 0.06, 2.68, size=19)
    x += w + gap
card(s, M, 3.85, 5.75, 2.2, fill=SOFT)
text(s, "KNOWLEDGE BASE", M + 0.32, 4.05, 5, 0.35, 11.5, ACC2, bold=True)
text(s, "The playbook is plain English. Every decision a lawyer makes is captured as a rule with the "
        "redline they used. It fills with PortSwigger's own positions, so the system knows more each week.",
     M + 0.32, 4.4, 5.15, 1.5, 13, INK2, spacing=1.2)
card(s, M + 6.05, 3.85, SW - M - (M + 6.05), 2.2, fill=SOFT)
text(s, "SKILLS THAT COMPOUND", M + 6.37, 4.05, 5, 0.35, 11.5, ACC2, bold=True)
text(s, "A new contract type or clause is not new code, it is a new skill added to the playbook. "
        "The agent's abilities grow without an engineer, and none of it is lost when I leave.",
     M + 6.37, 4.4, SW - M - (M + 6.05) - 0.6, 1.5, 13, INK2, spacing=1.2)
text(s, "Demonstrated: python evals/demo_learning.py closes the one gap the eval found, recall 93% to 100%, "
        "one playbook edit, no code.", M, 6.25, SW - 2 * M, 0.4, 12.5, INK2, italic=True)

# =============================================================================
# 13. PRIVATE BY DESIGN + STANDARDS
# =============================================================================
s = slide()
title(s, "Private by design", "Local AI on your server, and built to the standards that matter.")
card(s, M, 2.0, 5.55, 2.35, fill=SOFT)
text(s, "WHY LOCAL", M + 0.32, 2.18, 4, 0.35, 11.5, ACC2, bold=True)
text(s, "Contracts are privileged and confidential, and you are a security company. So the reasoning "
        "model runs on your own server, not a third-party API. No contract text ever leaves your "
        "infrastructure: privacy is kept by construction, not by policy.",
     M + 0.32, 2.52, 5.0, 1.7, 13, INK2, spacing=1.22)
card(s, M + 5.9, 2.0, SW - M - (M + 5.9), 2.35, fill=SOFT)
text(s, "HOW", M + 6.22, 2.18, 4, 0.35, 11.5, ACC2, bold=True)
text(s, "The deterministic engine needs no model at all: that is the floor. The reasoning layer is "
        "an open-weight model (Hermes 3, Llama 3.3, Qwen, Mistral) served locally with Ollama or "
        "vLLM. It ships as a skill on day one, then moves fully in-house.",
     M + 6.22, 2.52, SW - M - (M + 5.9) - 0.6, 1.7, 13, INK2, spacing=1.22)
text(s, "BUILT TO THE STANDARDS THAT MATTER", M, 4.62, 8, 0.35, 11.5, ACC2, bold=True)
std = [("UK GDPR / DPA 2018", "Personal data in contracts, handled lawfully"),
       ("ISO/IEC 27001", "Information security, on-prem by default"),
       ("ISO/IEC 42001", "AI management system, governed and audited"),
       ("EU AI Act / NIST AI RMF", "Human-in-the-loop, explainable, risk-based")]
sw2 = (SW - 2 * M - 3 * 0.24) / 4
for i, (name, sub) in enumerate(std):
    x = M + i * (sw2 + 0.24)
    card(s, x, 4.98, sw2, 1.35, fill=WHITE, line=LINE)
    text(s, name, x + 0.22, 5.14, sw2 - 0.36, 0.6, 12.5, INK, bold=True, spacing=1.05)
    text(s, sub, x + 0.22, 5.66, sw2 - 0.38, 0.6, 10.8, DIM, spacing=1.1)

# =============================================================================
# 14. WHY I APPLIED + HOW I ADD VALUE
# =============================================================================
s = slide()
title(s, "My fit", "Why I applied, and how I would add value here.")
text(s, "WHY I APPLIED", M, 2.05, 5, 0.35, 12, ACC2, bold=True)
multi(s, [
    ("PortSwigger is treating AI as the way it works, not a side project.", INK, True),
    ("I already do this: drop into a domain, build a real thing, measure it, hand it over.", DIM, False),
    ("I want to be somewhere the ambition and the tooling are taken seriously.", DIM, False),
], M, 2.45, 5.35, 2.6, 13.5, spacing=1.25, space_after=10)
text(s, "HOW I ADD VALUE", M + 6.0, 2.05, 5, 0.35, 12, ACC2, bold=True)
vals = [("Ship artefacts, not slides", "A working tool someone uses on Monday."),
        ("Bring the expert with me", "I work alongside the domain expert, not above them."),
        ("Measure everything", "I care as much about the eval as the build."),
        ("Leave capability, not dependency", "The playbook and skills stay when I move on.")]
y = 2.45
for hd, sub in vals:
    circle(s, "", M + 6.0, y + 0.03, 0.2, fill=ACC)
    text(s, hd, M + 6.4, y - 0.04, SW - M - (M + 6.4), 0.35, 13.5, INK, bold=True)
    text(s, sub, M + 6.4, y + 0.3, SW - M - (M + 6.4), 0.35, 12, DIM)
    y += 0.82

# =============================================================================
# 15. CLOSE
# =============================================================================
s = slide(footer=False)
dot(s, GREEN, M, 1.5, 0.2); dot(s, AMBER, M + 0.32, 1.5, 0.2); dot(s, RED, M + 0.64, 1.5, 0.2)
text(s, "Thank you.", M - 0.03, 1.95, 11, 0.9, 52, INK, bold=True)
text(s, "This was one corner. There are twenty-five more.", M, 3.1, 11.5, 0.6, 24, ACC2, bold=True)
multi(s, [
    ("The legal queue was the brief. The method is the point: find the risk gate, insert triage, "
     "measure it, keep it private, teach it to compound, then move on.", INK, False),
    ("Legal today. Support, marketing, finance, recruitment next. Same loop, new domain.", DIM, False),
], M, 4.0, 10.9, 1.5, 16, spacing=1.25, space_after=12)
text(s, "Questions?", M, 5.75, 6, 0.5, 20, INK, bold=True)
text(s, "Abdelrhman Rayis   ·   abdelrhman.rayis@gmail.com", M, 6.35, 9, 0.4, 13, DIM)
text(s, "Prototype, live console, evals and this deck are all in the submitted repo.", M, 6.7, 9, 0.4, 11.5, FAINT)

# ---- save --------------------------------------------------------------------
out = ROOT / "deck" / "FirstPass_PortSwigger.pptx"
prs.save(str(out))
print(f"Saved {out}  ·  {len(prs.slides._sldIdLst)} slides")
